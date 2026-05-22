"""
PowerPoint Pathway Diagram Composer
====================================
Assemble molecular pathway diagrams as editable PowerPoint (.pptx) files.
Use when manual tweaking of node positions, arrow paths, or labels is needed
before final export to PDF/PNG.

Usage:
    python -m smiles_to_3d.pathway_pptx spec.yaml --img-dir ./pngs/ --output pathway.pptx

Dependencies:
    pip install python-pptx Pillow PyYAML lxml
"""

import argparse
import os
import math
import yaml
from collections import defaultdict
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Defaults ──────────────────────────────────────────────────────────────

SLIDE_W_IN = 13.33   # 16:9
SLIDE_H_IN = 7.5
MAIN_H_IN = 1.0
SIDE_H_IN = 0.65
MAIN_Y_IN = 2.2
SIDE_Y_IN = 0.9
MARGIN_X_IN = 0.6
ARROW_GAP = 0.35

RGB_SOLID = RGBColor(0x2C, 0x2C, 0x2C)
RGB_DASHED = RGBColor(0xAA, 0xAA, 0xAA)
RGB_TEXT = RGBColor(0x11, 0x11, 0x11)
RGB_SUBTEXT = RGBColor(0x55, 0x55, 0x55)
RGB_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

NSMAP_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
NSMAP_P = "http://schemas.openxmlformats.org/presentationml/2006/main"


def _emu(inches_val):
    return int(inches_val * 914400)


def _color_hex(color):
    return f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"


def _add_line(slide, x1_in, y1_in, x2_in, y2_in, color, width_pt,
              arrow_end=True, dashed=False):
    """Add a connector line with optional arrowhead via raw XML."""
    x1 = _emu(x1_in) if isinstance(x1_in, (int, float)) else int(x1_in)
    y1 = _emu(y1_in) if isinstance(y1_in, (int, float)) else int(y1_in)
    x2 = _emu(x2_in) if isinstance(x2_in, (int, float)) else int(x2_in)
    y2 = _emu(y2_in) if isinstance(y2_in, (int, float)) else int(y2_in)

    w_emu = int(Pt(width_pt) * 12700)
    dash_xml = '<a:prstDash val="dash"/>' if dashed else ""
    tail_xml = '<a:tailEnd type="triangle" w="med" len="med"/>' if arrow_end else ""

    dx = x2 - x1
    dy = y2 - y1

    cxn = etree.fromstring(f"""
    <p:cxnSp xmlns:p="{NSMAP_P}" xmlns:a="{NSMAP_A}">
      <p:nvCxnSpPr>
        <p:cNvPr id="0" name=""/>
        <p:cNvCxnSpPr/>
        <p:nvPr/>
      </p:nvCxnSpPr>
      <p:spPr>
        <a:xfrm>
          <a:off x="{x1}" y="{y1}"/>
          <a:ext cx="{dx}" cy="{dy}"/>
        </a:xfrm>
        <a:prstGeom prst="line"><a:avLst/></a:prstGeom>
        <a:ln w="{w_emu}">
          <a:solidFill><a:srgbClr val="{_color_hex(color)}"/></a:solidFill>
          {dash_xml}
          {tail_xml}
        </a:ln>
      </p:spPr>
    </p:cxnSp>
    """)

    slide.shapes._spTree.append(cxn)


def _add_label(slide, text, cx, y, font_size_pt, color, bold=True, italic=False):
    """Add a centered text label."""
    tf = slide.shapes.add_textbox(
        Inches(cx - 0.6), Inches(y), Inches(1.2), Inches(0.3)
    )
    tf.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.text_frame.paragraphs[0].add_run()
    run.text = text
    run.font.name = "Times New Roman"
    run.font.size = Pt(font_size_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tf


def compute_layout(spec, img_dir):
    """Compute node positions in inches."""
    main_chain = spec["main_chain"]
    side_fragments = spec.get("side_fragments", [])

    sample = os.path.join(img_dir, main_chain[0]["png"])
    orig_w, orig_h = Image.open(sample).size

    ratio = MAIN_H_IN / orig_h
    main_w = orig_w * ratio
    side_w = orig_w * (SIDE_H_IN / orig_h)

    n = len(main_chain)
    min_arrow = 0.6
    gap_x = max(main_w + ARROW_GAP * 2 + 0.5, main_w + min_arrow + 0.3)

    # Group side fragments
    groups = defaultdict(list)
    for sf in side_fragments:
        groups[sf.get("attach_to", 0)].append(sf)

    # Auto-scaling if content exceeds slide width
    available_w = SLIDE_W_IN - MARGIN_X_IN * 2
    total_content_w = main_w * n + (n - 1) * gap_x

    if total_content_w > available_w:
        scale = available_w / total_content_w * 0.95
        main_w *= scale
        side_w *= scale
        gap_x *= scale

    # Main chain nodes
    nodes = []
    for i, node in enumerate(main_chain):
        x = MARGIN_X_IN + i * (main_w + gap_x)
        cx = x + main_w / 2
        cy = MAIN_Y_IN + MAIN_H_IN / 2
        nodes.append({
            "id": node["id"],
            "png": node["png"],
            "label": node.get("label", ""),
            "label_cn": node.get("label_cn", ""),
            "arrow_label": node.get("arrow_label", ""),
            "role": "main",
            "x": x, "y": MAIN_Y_IN,
            "w": main_w, "h": MAIN_H_IN,
            "cx": cx, "cy": cy,
            "node_idx": i,
        })

    # Side fragments
    side_nodes = []
    side_arrows = []

    for attach_idx, fragments in groups.items():
        if attach_idx >= n:
            continue
        parent = nodes[attach_idx]
        for j, sf in enumerate(fragments):
            ox = parent["x"] + parent["w"] + ARROW_GAP + 0.15 + j * (side_w + 0.3)
            oy = SIDE_Y_IN + SIDE_H_IN / 2
            scx = ox + side_w / 2

            side_nodes.append({
                "id": sf["id"],
                "png": sf["png"],
                "label": sf.get("label", ""),
                "role": "side",
                "x": ox, "y": oy,
                "w": side_w, "h": SIDE_H_IN,
                "cx": scx, "cy": oy,
                "attach_to": attach_idx,
            })

            side_arrows.append({
                "type": "dashed",
                "sx": parent["x"] + parent["w"],
                "sy": parent["cy"],
                "tx": ox,
                "ty": oy,
                "label": sf.get("arrow_label", ""),
            })

    # Main chain arrows
    main_arrows = []
    for i in range(n - 1):
        L = nodes[i]
        R = nodes[i + 1]
        main_arrows.append({
            "type": "solid",
            "sx": L["x"] + L["w"],
            "sy": L["cy"],
            "tx": R["x"],
            "ty": R["cy"],
            "label": L.get("arrow_label", ""),
        })

    return nodes + side_nodes, main_arrows + side_arrows


def render_pptx(spec, img_dir, output_path):
    """Render a pathway diagram as PPTX."""
    main_chain = spec["main_chain"]
    side_fragments = spec.get("side_fragments", [])

    nodes, arrows = compute_layout(spec, img_dir)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # White background
    bg = slide.shapes.add_shape(1, 0, 0, Inches(SLIDE_W_IN), Inches(SLIDE_H_IN))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGB_WHITE
    bg.line.fill.background()

    # Place molecule images
    for node in nodes:
        img_path = os.path.join(img_dir, node["png"])
        if not os.path.exists(img_path):
            print(f"  ⚠ Missing: {img_path}")
            continue
        slide.shapes.add_picture(
            img_path,
            Inches(node["x"]), Inches(node["y"]),
            Inches(node["w"]), Inches(node["h"]),
        )

    # Draw arrows
    for arrow in arrows:
        if arrow["type"] == "solid":
            _add_line(slide,
                Inches(arrow["sx"]), Inches(arrow["sy"]),
                Inches(arrow["tx"]), Inches(arrow["ty"]),
                RGB_SOLID, 1.6, arrow_end=True, dashed=False)
            lbl = arrow.get("label", "")
            if lbl:
                cx = (arrow["sx"] + arrow["tx"]) / 2
                _add_label(slide, lbl, cx, arrow["sy"] - 0.18,
                    11, RGBColor(0x44, 0x44, 0x44), bold=False, italic=False)
        else:
            # Dashed S-curve: three-segment approximation
            mid_x = (arrow["sx"] + arrow["tx"]) / 2
            _add_line(slide,
                Inches(arrow["sx"]), Inches(arrow["sy"]),
                Inches(mid_x), Inches(arrow["sy"]),
                RGB_DASHED, 1.3, arrow_end=False, dashed=True)
            _add_line(slide,
                Inches(mid_x), Inches(arrow["sy"]),
                Inches(mid_x), Inches(arrow["ty"]),
                RGB_DASHED, 1.3, arrow_end=False, dashed=True)
            _add_line(slide,
                Inches(mid_x), Inches(arrow["ty"]),
                Inches(arrow["tx"]), Inches(arrow["ty"]),
                RGB_DASHED, 1.3, arrow_end=True, dashed=True)

    # Labels
    for node in nodes:
        cx = node["cx"]
        if node["role"] == "main":
            _add_label(slide, node["label"], cx, node["y"] - 0.10,
                12, RGB_TEXT, bold=True, italic=False)
            if node.get("label_cn"):
                _add_label(slide, node["label_cn"], cx, node["y"] - 0.28,
                    9, RGB_SUBTEXT, bold=False, italic=True)
        else:
            _add_label(slide, node["label"], cx, node["y"] + node["h"] + 0.04,
                10, RGB_TEXT, bold=True, italic=False)

    prs.save(output_path)
    print(f"✓ Saved: {output_path}")
    return output_path


# ── CLI ────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="Compose molecular pathway diagrams as PPTX"
    )
    parser.add_argument("spec", help="YAML topology spec file")
    parser.add_argument("--img-dir", "-i", required=True,
                        help="Directory containing molecular PNG images")
    parser.add_argument("--output", "-o", required=True,
                        help="Output PPTX file")
    args = parser.parse_args()

    with open(args.spec, encoding="utf-8") as f:
        spec = yaml.safe_load(f)

    render_pptx(spec, args.img_dir, args.output)


if __name__ == "__main__":
    main()
